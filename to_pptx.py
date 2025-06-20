from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


# --- Helper function to add a content slide ---
def add_content_slide(prs, title_text, content_list, notes=""):
    layout = prs.slide_layouts[5]  # Choose a layout (e.g., Title and Content)
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    title.text = title_text

    body = slide.placeholders[1]  # Body placeholder
    tf = body.text_frame
    tf.clear()  # Clear existing text

    for item in content_list:
        p = tf.add_paragraph()
        if isinstance(item, tuple) and item[0] == "code":
            run = p.add_run()
            run.text = item[1]
            run.font.name = 'Courier New'
            run.font.size = Pt(10)  # Smaller for code
        elif isinstance(item, tuple) and item[0] == "bold":
            run = p.add_run()
            run.text = item[1]
            run.font.bold = True
        elif isinstance(item, tuple) and item[0] == "bullet":
            p.text = item[1]
            p.level = item[2] if len(item) > 2 else 0
        else:
            p.text = item

    if notes:  # Add notes to the notes slide
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes


# --- Create Presentation ---
prs = Presentation()

# Slide 1: Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Database Management Systems (DSC 12)"
subtitle.text = "SQL Practicals (Q1-52)\nPresented by: Abitatha Roy (152304)"

# Slide 2: Deduced Database Schema
schema_content = [
    ("bold", "Deduced Database Schema"),
    "Based on the queries, we're likely working with the following tables:",
    "",
    ("bold", "Table: `employee`"),
    ("bullet", "Eno (INT/VARCHAR): Employee Number (Primary Key)", 1),
    ("bullet", "Ename (VARCHAR): Employee Name", 1),
    ("bullet", "Job_type (VARCHAR): Job Title", 1),
    ("bullet", "Manager (INT/VARCHAR): Manager's Employee Number (FK to employee.Eno, Can be NULL)", 1),
    ("bullet", "Hire_date (DATE): Hiring Date", 1),
    ("bullet", "Dno (INT): Department Number (FK to department.Dno)", 1),
    ("bullet", "Commission (DECIMAL): Commission earned (Can be NULL)", 1),
    ("bullet", "Salary (DECIMAL): Monthly Salary", 1),
    "",
    ("bold", "Table: `department`"),
    ("bullet", "Dno (INT): Department Number (Primary Key)", 1),
    ("bullet", "Dname (VARCHAR): Department Name", 1),
    ("bullet", "Location (VARCHAR): Department Location", 1),
]
add_content_slide(prs, "Understanding Our Database Schema", schema_content,
                  "Verbally elaborate on data types, potential constraints, and relationships.")

# --- Query Data (Abbreviated for brevity in this example, you'll fill all 52) ---
queries_data = [
    {
        "num": 1,
        "title_short": "Basic Employee Details",
        "objective": "Query to display Employee Name, Job, Hire Date, Employee Number; for each employee with the Employee Number appearing first.",
        "sql": """SELECT Eno, Ename, Job_type, Hire_date
FROM employee;""",
        "notes": ""
    },
    {
        "num": 2,
        "title_short": "Unique Job Types",
        "objective": "Query to display Unique Jobs from the Employee Table.",
        "sql": """SELECT DISTINCT(Job_type)
FROM employee;""",
        "notes": ""
    },
    {
        "num": 3,
        "title_short": "Concatenating Name and Job",
        "objective": "Query to display the Employee Name concatenated by a Job separated by a comma.",
        "sql": """SELECT CONCAT(Ename, ", ", Job_type) AS Employee_Job
FROM employee;""",
        "notes": "Added an alias `Employee_Job` for clarity."
    },
    {
        "num": 4,
        "title_short": "All Data Concatenated",
        "objective": "Query to display all the data from the Employee Table. Separate each Column by a comma and name the said column as THE_OUTPUT.",
        "sql": """SELECT CONCAT(Eno, ", ", Ename, ", ", Job_type, ", ", Manager, ", ", Hire_date, ", ", Dno, ", ", Commission, ", ", Salary) AS THE_OUTPUT
FROM employee;""",
        "notes": ""
    },
    {
        "num": 5,
        "title_short": "High Earners",
        "objective": "Query to display the Employee Name & Salary of all the employees earning more than $2850.",
        "sql": """SELECT Ename, Salary
FROM employee
WHERE Salary > 2850;""",
        "notes": ""
    },
    {
        "num": 6,
        "title_short": "Specific Employee",
        "objective": "Query to display Employee Name & Department Number for the Employee No= 7900.",
        "sql": """SELECT Ename, Dno
FROM employee
WHERE Eno = 7900;""",
        "notes": ""
    },
    {
        "num": 7,
        "title_short": "Salary Range Exclusion",
        "objective": "Query to display Employee Name & Salary for all employees whose salary is not in the range of $1500 and $2850.",
        "sql": """SELECT Ename, Salary
FROM employee
WHERE Salary NOT BETWEEN 1500 AND 2850;""",
        "notes": ""
    },
    {
        "num": 8,
        "title_short": "Hires in Period",
        "objective": "Query to display Employee Name, Job, and Hire Date of all the employees hired between Feb 20, 1981 and May 1, 1981. Order the query in ascending order of Start Date.",
        "sql": """SELECT Ename, Job_type, Hire_date
FROM employee
WHERE Hire_date BETWEEN '1981-02-20' AND '1981-05-01'
ORDER BY Hire_date ASC;""",
        "notes": ""
    },
    {
        "num": 9,
        "title_short": "Specific Departments",
        "objective": "Query to display Employee Name & Department No. of all the employees in Dept 10 and Dept 30 in the alphabetical order by name.",
        "sql": """SELECT Ename, Dno
FROM employee
WHERE Dno = 10 OR Dno = 30 -- Alternative: Dno IN (10, 30)
ORDER BY Ename ASC;""",
        "notes": ""
    },
    {
        "num": 10,
        "title_short": "High Earners in Depts",
        "objective": "Query to display Employee Name & Salary of employees who earned more than $1500 and are in Department 10 or 30.",
        "sql": """SELECT Ename, Salary
FROM employee
WHERE (Dno = 10 OR Dno = 30) AND
	Salary > 1500;""",
        "notes": ""
    },
    {
        "num": 11,
        "title_short": "Hired in 1981",
        "objective": "Query to display Name & Hire Date of every Employee who was hired in 1981.",
        "sql": """SELECT Ename, Hire_date
FROM employee
WHERE YEAR(Hire_date) = 1981; -- Or '1981'""",
        "notes": "YEAR() extracts year. Numeric comparison is robust."
    },
    {
        "num": 12,
        "title_short": "No Manager",
        "objective": "Query to display Name & Job of all employees who don’t have a current Manager.",
        "sql": """SELECT Ename, Job_type
FROM employee
WHERE Manager IS NULL;""",
        "notes": ""
    },
    {
        "num": 13,
        "title_short": "Commission Earners (Sorted)",
        "objective": "Query to display the Name, Salary & Commission for all the employees who earn commission. Sort the data in descending order of Salary and Commission.",
        "sql": """SELECT Ename, Salary, Commission
FROM employee
WHERE Commission IS NOT NULL AND Commission > 0 -- Original: Commission > 0
ORDER BY Salary DESC, Commission DESC; -- Original: ASC, ASC""",
        "notes": "Adjusted sort order to DESC as per objective text. Original query used ASC."
    },
    {
        "num": 14,
        "title_short": "Name Pattern (3rd 'A')",
        "objective": "Query to display Name of all the employees where the third letter of their name is ‘A’.",
        "sql": """SELECT Ename
FROM employee
WHERE Ename LIKE '__A%';""",
        "notes": ""
    },
    {
        "num": 15,
        "title_short": "Complex Conditions",
        "objective": "Query to display Name of all employees either have two ‘R’s or have two ‘A’s in their name & are either in Dept No = 30 or their Manager’s Employee No = 7788.",
        "sql": """SELECT Ename
FROM employee
WHERE (Ename LIKE '%A%A%' OR  Ename LIKE '%R%R%') AND 
	(Dno = 30 OR Manager = 7788);""",
        "notes": ""
    },
    {
        "num": 16,
        "title_short": "Job & Salary Exclusions",
        "objective": "Query to display Name, Job and Salary of all employees whose Job is Clerical or Analyst & their salaries are not equal to 1000, 3000, or 5000.",
        "sql": """SELECT Ename, Job_type, Salary
FROM employee
WHERE (Job_type = 'Clerk' OR Job_type = 'Analyst') AND
		(Salary NOT IN (1000, 3000, 5000));""",
        "notes": ""
    },
    {
        "num": 17,
        "title_short": "Commission vs Salary",
        "objective": "Query to display Name, Salary and Commission for all employees whose Commission Amount is greater than their Salary increased by 5 %.",
        "sql": """SELECT Ename, Salary, Commission
FROM employee
WHERE Commission > 1.05 * Salary;""",
        "notes": ""
    },
    {
        "num": 18,
        "title_short": "Current Date",
        "objective": "Query to display the Current Date.",
        "sql": """SELECT CURDATE();""",
        "notes": ""
    },
    {
        "num": 19,
        "title_short": "Salary Increase (15%)",
        "objective": "Query to display Employee No., Name, Salary and the Salary increased by 15 % expressed as a absolute whole number.",
        "sql": """SELECT Eno, Ename, Salary, ROUND(Salary * 1.15) AS Increased_Salary
FROM employee;""",
        "notes": ""
    },
    {
        "num": 20,
        "title_short": "Salary Review Date",
        "objective": "Query to display Name, Hire Date and Salary Review Date which is the 1st Monday after six months of employment.",
        "sql": """SELECT Ename, 
		Hire_date,
        DATE_ADD(Hire_date, INTERVAL 6 MONTH) + INTERVAL((8 - DAYOFWEEK(DATE_ADD(Hire_date, INTERVAL 6 MONTH))) % 7) DAY
        AS Salary_Review_Date
FROM employee;""",
        "notes": "This date logic is specific. DAYOFWEEK() behavior can vary."
    },
    {
        "num": 21,
        "title_short": "Salary > Any Clerk",
        "objective": "Query to display the employees that earn a salary that is higher than the salary of any of the clerks.",
        "sql": """SELECT Eno, Ename
FROM employee
WHERE Salary > ANY (SELECT Salary
				FROM employee
                WHERE Job_type = 'Clerk');""",
        "notes": "'> ANY' means greater than the minimum of the subquery results."
    },
    {
        "num": 22,
        "title_short": "Months of Service",
        "objective": "Query to display Name and calculate the number of months between today and the date each employee was hired.",
        "sql": """SELECT Ename, TIMESTAMPDIFF(MONTH, Hire_date, CURDATE()) AS Date_Difference
FROM employee;""",
        "notes": "TIMESTAMPDIFF is MySQL specific."
    },
    {
        "num": 23,
        "title_short": "Dream Salary String",
        "objective": "Query to display for each employee: <E-Name> earns <Salary> monthly but wants < 3 * Current Salary >. Label the Column as Dream Salary.",
        "sql": """SELECT CONCAT(Ename, " earns ", Salary, " monthly but wants ", 3 * Salary) AS Dream_Salary
FROM employee;""",
        "notes": ""
    },
    {
        "num": 24,
        "title_short": "Formatted Salary (LPAD)",
        "objective": "Query to display Name and Salary for all employees. Format the salary to be 15 character long, left padded with $ sign.",
        "sql": """SELECT Ename, LPAD(Salary, 15, '$') AS Formatted_Salary
FROM employee;""",
        "notes": ""
    },
    {
        "num": 25,
        "title_short": "Name Format & Length",
        "objective": "Query to display Name with the 1st letter capitalized and all other letter lower case & length of their name of all the employees whose name starts with ‘J’,’A’ and ‘M’.",
        "sql": """SELECT CONCAT(UPPER(LEFT(Ename, 1)), LOWER(SUBSTRING(Ename, 2))) AS Formatted_Name, 
		LENGTH(Ename) AS Name_Length
FROM employee
WHERE Ename LIKE 'J%' OR 
		Ename LIKE 'A%' OR 
        Ename LIKE 'M%';""",
        "notes": ""
    },
    {
        "num": 26,
        "title_short": "Day of Week Hired",
        "objective": "Query to display Name, Hire Date and Day of the week on which the employee started.",
        "sql": """SELECT Ename, Hire_date, DAYNAME(Hire_date) AS 'Day when employee started'
FROM employee;""",
        "notes": ""
    },
    {
        "num": 27,
        "title_short": "Commission Status (CASE)",
        "objective": "Query to display Name and Commission Amount. If the employee does not earn commission then use default value ‘No Commission’.",
        "sql": """SELECT Ename,
       CASE 
           WHEN Commission IS NULL OR Commission = 0 THEN 'No Commission'
           ELSE CAST(Commission AS CHAR) 
       END AS Commission_Status
FROM employee;""",
        "notes": ""
    },
    {
        "num": 28,
        "title_short": "Employee & Dept (JOIN)",
        "objective": "Query to display Name, Department Name and Department No for all the employees.",
        "sql": """SELECT e.Ename, d.Dname, d.Dno
FROM employee e 
INNER JOIN department d
ON e.Dno = d.Dno;""",
        "notes": "First explicit JOIN with the department table."
    },
    {
        "num": 29,
        "title_short": "Unique Jobs in Dept 30",
        "objective": "Query to display Unique Listing of all Jobs that are in Department #30.",
        "sql": """SELECT DISTINCT(Job_type)
FROM employee
WHERE Dno = 30;""",
        "notes": ""
    },
    {
        "num": 30,
        "title_short": "Commission Earners & Dept",
        "objective": "Query to display Name, Department Name and Location for all employees earning a commission.",
        "sql": """SELECT e.Ename, d.Dname, d.Location
FROM employee e
INNER JOIN department d 
ON e.Dno = d.Dno 
WHERE Commission IS NOT NULL AND Commission > 0;""",
        "notes": ""
    },
    {
        "num": 31,
        "title_short": "'A' in Name & Dept",
        "objective": "Query to display Name, Dept Name of all employees who have an ‘A’ in their name.",
        "sql": """SELECT e.Ename, d.Dname
FROM employee e
INNER JOIN department d 
ON e.Dno = d.Dno 
WHERE e.Ename LIKE "%A%";""",
        "notes": ""
    },
    {
        "num": 32,
        "title_short": "Employees in 'Dallas'",
        "objective": "Query to display Name, Job, Department No. and Department Name for all the employees working at the Dallas location.",
        "sql": """SELECT e.Ename, e.Job_type, e.Dno, d.Dname 
FROM employee e 
INNER JOIN department d
ON e.Dno = d.Dno 
WHERE d.Location = "Dallas";""",
        "notes": ""
    },
    {
        "num": 33,
        "title_short": "Employee & Manager (Self JOIN)",
        "objective": "Query to display Name and Employee No. along with their Manager’s Name and Manager’s employee no.",
        "sql": """SELECT e.Ename AS "Employee Name", 
		e.Eno AS "Employee Number", 
        m.Ename AS "Manager Name", 		
        m.Eno AS "Manager Eno" -- Corrected from original text
FROM employee e 
JOIN employee m 	
ON e.Manager = m.Eno;""",
        "notes": "Self JOIN. Original had an anomaly, corrected to typical self-join for manager."
    },
    {
        "num": 34,
        "title_short": "Emp/Mgr (LEFT JOIN)",
        "objective": "Query to display Name and Employee no. along with their Manger’s Name and the Manager’s employee no; along with the Employees’ Name who do not have a Manager.",
        "sql": """SELECT e.Ename AS "Employee Name", 
		e.Eno AS "Employee Number", 
        COALESCE(m.Ename, "NA") AS "Manager Name",
        COALESCE(CAST(m.Eno AS CHAR), "NA") AS "Manager Eno" -- Cast for COALESCE with string
FROM employee e 
LEFT JOIN employee m
ON e.Manager = m.Eno;""",
        "notes": "LEFT JOIN to include employees without managers. COALESCE handles NULLs."
    },
    {
        "num": 35,
        "title_short": "Salary > Avg & 'T' in Dept",
        "objective": "Query to display the Employee No, Name and Salary for all employees who earn more than the average salary and who work in a Department with any employee with a ‘T’ in his/her name.",
        "sql": """SELECT e.Eno, e.Ename, e.Salary 
FROM employee e 
WHERE e.Salary > (SELECT AVG(Salary) FROM employee) AND
		e.Dno IN (
			SELECT DISTINCT Dno -- Use DISTINCT
            FROM employee
            WHERE Ename LIKE "%T%" -- Original: "%T"
		);""",
        "notes": ""
    },
    {
        "num": 36,
        "title_short": "Matching Dept & Salary",
        "objective": "Query to display Name, Dept No. & Salary of any employee whose department No. and salary matches both the department no. and the salary of any employee who earns a commission.",
        "sql": """SELECT e.Ename, e.Dno, e.Salary
FROM employee e 
WHERE (e.Dno, e.Salary) IN (SELECT f.Dno, f.Salary
				FROM employee f
                WHERE f.Commission != 0 AND f.Commission IS NOT NULL
                );""",
        "notes": "Tuple comparison in subquery."
    },
    {
        "num": 37,
        "title_short": "Hired After 'Blake'",
        "objective": "Query to display Name, Hire Date of any employee hired after the employee Blake was hired by the Company.",
        "sql": """SELECT Ename, Hire_date
FROM employee
WHERE Hire_date > (SELECT Hire_date
					FROM employee
                    WHERE Ename = "Blake"
                    );""",
        "notes": ""
    },
    {
        "num": 38,
        "title_short": "Hired Before Manager",
        "objective": "Query to display Name and Hire Dates of all Employees along with their Manager’s Name and Hire Date for all the employees who were hired before their managers.",
        "sql": """SELECT e.Ename AS "Employee Name", e.Hire_date AS "Employee Hire Date",
       m.Ename AS "Manager Name", m.Hire_date AS "Manager Hire Date"
FROM employee e
INNER JOIN employee m
ON e.Manager = m.Eno
WHERE e.Hire_date < m.Hire_date;""",
        "notes": ""
    },
    {
        "num": 39,
        "title_short": "Salary as Asterisks",
        "objective": "Query to display Name and Salaries represented by Asteristisks – “Each asterisks (*) signifying $100.",
        "sql": """SELECT Ename,
		REPEAT('*', FLOOR(Salary/100)) AS Salary_In_Asterisks -- Use FLOOR
FROM employee;""",
        "notes": "REPEAT is MySQL specific. FLOOR for non-even hundreds."
    },
    {
        "num": 40,
        "title_short": "Overall Salary Stats",
        "objective": "Query to display the Highest, Lowest, Sum and Average Salaries of all the employees.",
        "sql": """SELECT 
    MAX(Salary) AS Max_Salary,
    MIN(Salary) AS Min_Salary,
    SUM(Salary) AS Sum_Salary,
    AVG(Salary) AS Avg_Salary
FROM employee;""",
        "notes": ""
    },
    {
        "num": 41,
        "title_short": "Stats by Job Type",
        "objective": "Query to display Highest, Lowest, Sum and Average Salary for each unique Job Type.",
        "sql": """SELECT 
	Job_type, -- DISTINCT removed due to GROUP BY
    MAX(Salary) AS Max_Salary,
    MIN(Salary) AS Min_Salary,
    SUM(Salary) AS Sum_Salary,
    AVG(Salary) AS Avg_Salary
FROM employee
GROUP BY Job_type;""",
        "notes": ""
    },
    {
        "num": 42,
        "title_short": "Count by Job Type",
        "objective": "Query to display the number of employees performing the same Job type functions.",
        "sql": """SELECT Job_type, -- DISTINCT removed
		COUNT(*) AS "No. of Employees" -- COUNT(Ename) is also fine
FROM employee
GROUP BY Job_type;""",
        "notes": ""
    },
    {
        "num": 43,
        "title_short": "Number of Managers",
        "objective": "Query to display the no. of managers without listing their names.",
        "sql": """SELECT COUNT(DISTINCT Manager) AS "Number of Managers"
FROM employee
WHERE Manager IS NOT NULL; -- Important if Manager can be NULL""",
        "notes": ""
    },
    {
        "num": 44,
        "title_short": "Salary Range Diff",
        "objective": "Query to display the Difference b/w the Highest and Lowest Salaries.",
        "sql": """SELECT (MAX(Salary) - MIN(Salary)) AS "Difference between Highest and Lowest Salary"
FROM employee;""",
        "notes": ""
    },
    {
        "num": 45,
        "title_short": "Min Salary per Mgr (HAVING)",
        "objective": "Query to display the Manager’s No. & the Salary of the Lowest paid employee for that respective manager. Exclude anyone where the Manager ID is not known. Exclude any groups where the minimum salary is less than $1000.",
        "sql": """SELECT 
	Manager AS Manager_Eno,
    MIN(Salary) AS "Min. Salary"
FROM employee
WHERE Manager IS NOT NULL
GROUP BY Manager
HAVING MIN(Salary) >= 1000; -- Keep if min salary is $1000 or more""",
        "notes": "HAVING filters groups after aggregation."
    },
    {
        "num": 46,
        "title_short": "Department Statistics",
        "objective": "Query to display the Department Name, Location Name, No. of Employees & the average salary for all employees in that department.",
        "sql": """SELECT 
	D.Dname,
    D.Location,
    COUNT(E.Eno) AS Number_Of_Employees,
    AVG(E.Salary) AS Avg_Salary_In_Dept
FROM department D
JOIN employee E ON D.Dno = E.Dno
GROUP BY D.Dno, D.Dname, D.Location;""",
        "notes": "Group by all non-aggregated selected columns from department."
    },
    {
        "num": 47,
        "title_short": "Employees in Blake's Dept",
        "objective": "Query to display Name and Hire Date for all employees in the same dept. as Blake.",
        "sql": """SELECT Ename, Hire_date
FROM employee
WHERE 
	Dno = (SELECT Dno
			FROM employee
            WHERE Ename = "Blake"
		);""",
        "notes": ""
    },
    {
        "num": 48,
        "title_short": "Salary > Average",
        "objective": "Query to display the Employee No. & Name for all employees who earn more than the average salary.",
        "sql": """SELECT Eno, Ename
FROM employee
WHERE Salary > (SELECT AVG(Salary)
				FROM employee
                );""",
        "notes": ""
    },
    {
        "num": 49,
        "title_short": "Dept with 'T'-Named Colleague",
        "objective": "Query to display Employee Number & Name for all employees who work in a department with any employee whose name contains a ‘T’.",
        "sql": """SELECT Eno, Ename
FROM employee
WHERE Dno IN (SELECT DISTINCT Dno -- Changed from = to IN
			FROM employee
            WHERE Ename LIKE '%T%'
            );""",
        "notes": "Changed to IN as subquery can return multiple departments."
    },
    {
        "num": 50,
        "title_short": "Reporting to 'King'",
        "objective": "Query to display the employee name and salary of all employees who report to King.",
        "sql": """SELECT Ename, Salary
FROM employee
WHERE Manager = (SELECT Eno
				FROM employee
                WHERE Ename = 'King'
                );""",
        "notes": ""
    },
    {
        "num": 51,
        "title_short": "Employees in 'Sales' Dept",
        "objective": "Query to display the Department No, Name & Job for all employees in the Sales Dept.",
        "sql": """SELECT e.Dno, e.Ename, e.Job_type
FROM employee e
JOIN department d ON e.Dno = d.Dno
WHERE d.Dname = "Sales";""",
        "notes": ""
    },
    {
        "num": 52,
        "title_short": "Manager Salary > Dept Avg",
        "objective": "Query to select manager name getting salary greater than average salary of employees in his department.",
        "sql": """SELECT M.Ename
FROM employee M
WHERE M.Eno IN (SELECT DISTINCT Manager FROM employee WHERE Manager IS NOT NULL) -- Ensures M is a manager
		AND
        M.Salary > (SELECT AVG(Salary)
					FROM employee
                    WHERE Dno = M.Dno -- Correlated subquery
                    );""",
        "notes": "Correlated subquery: inner query depends on outer query's current row."
    }
]

# Add slides for each query
for query in queries_data:
    content = [
        ("bold", "Objective:"),
        query["objective"],
        "",
        ("bold", "SQL Query:"),
        ("code", query["sql"])
    ]
    slide_title = f"Query {query['num']}: {query['title_short']}"
    add_content_slide(prs, slide_title, content, query.get("notes", ""))

# Slide: Conclusion
conclusion_content = [
    ("bold", "We've covered a wide range of SQL queries, including:"),
    ("bullet", "Basic Selections & Filtering (`SELECT`, `WHERE`, `ORDER BY`)", 1),
    ("bullet", "String, Date, and Numeric Functions", 1),
    ("bullet", "Aggregate Functions (`COUNT`, `SUM`, `AVG`, `MAX`, `MIN`)", 1),
    ("bullet", "Grouping Data (`GROUP BY`, `HAVING`)", 1),
    ("bullet", "Joining Tables (`INNER JOIN`, `LEFT JOIN`, Self Joins)", 1),
    ("bullet", "Subqueries (Scalar, Multi-row, Correlated)", 1),
    "",
    "These examples demonstrate how SQL can be used to retrieve, manipulate, and analyze data effectively.",
    "",
    ("bold", "Any Questions?")
]
add_content_slide(prs, "Conclusion & Q&A", conclusion_content)

# Save the presentation
file_path = "SQL_Lecture_Slides.pptx"
prs.save(file_path)
print(f"Presentation saved as {file_path}")